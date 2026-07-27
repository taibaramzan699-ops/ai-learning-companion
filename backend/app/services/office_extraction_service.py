"""Text extraction for Office Open XML formats (.docx, .pptx).

Mirrors the (page_number, text) tuple-list shape that ocr_service's PDF
functions return, so document_processing.py can feed either straight into
the same chunk_text() call without branching downstream.

Requires: python-docx, python-pptx (both pure-Python, no system deps —
add to requirements.txt: python-docx>=1.1.0, python-pptx>=0.6.23)
"""
from io import BytesIO

from docx import Document as DocxDocument
from pptx import Presentation


def extract_text_from_docx(file_bytes: bytes) -> list[tuple[int, str]]:
    """.docx has no fixed pagination (page breaks depend on the reader's
    layout engine, not the file format), so the whole document is returned
    as a single logical page. page_count will show as 1 for these documents —
    that's expected, not a bug."""
    doc = DocxDocument(BytesIO(file_bytes))

    parts: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]

    # Paragraphs alone miss table content — pull that in too.
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    parts.append(cell.text)

    full_text = "\n\n".join(parts)
    return [(1, full_text)]


def extract_text_from_pptx(file_bytes: bytes) -> list[tuple[int, str]]:
    """One slide = one 'page', which maps naturally onto the existing
    page_number concept and gives a meaningful page_count in the UI."""
    prs = Presentation(BytesIO(file_bytes))
    pages: list[tuple[int, str]] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs)
                    if line.strip():
                        lines.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            lines.append(cell.text)
        pages.append((slide_index, "\n".join(lines)))

    return pages